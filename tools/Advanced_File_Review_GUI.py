"""
Advanced File Review GUI
=========================
Supports: PDF, images (JPG/PNG/TIFF/BMP/GIF), Word (.docx),
          Excel (.xlsx), PowerPoint (.pptx)

Point this tool at any folder to browse its documents, add dated
notes per file, and export all notes to CSV or TXT.

Layout:
  [ Folder path bar                                              ]
  [ File List  |  Document Viewer (zoom / page nav)  |  Notes  ]

--------------------------------------------------------------
HOW TO RUN (Windows)
--------------------------------------------------------------
Do NOT double-click this file — it won't work that way.
Instead, launch it through PowerShell:

  1. Open File Explorer and navigate to the folder where
     this script lives.

  2. Click the address bar, type  powershell  and press Enter.
     (Use PowerShell, not CMD — especially on network drives \\)

  3. In the PowerShell window, type:
         python Advanced_File_Review_GUI.py
     and press Enter.

--------------------------------------------------------------
IF YOU GET "Missing dependency" OR AN IMPORT ERROR
--------------------------------------------------------------
Run this in PowerShell, then try again:

    python -m pip install pymupdf Pillow python-docx openpyxl python-pptx

python-docx, openpyxl, and python-pptx are optional — the viewer
will still open but will show an install prompt for those file types.

--------------------------------------------------------------
"""

import os
import sys
import json
import csv
import textwrap
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from datetime import datetime
from pathlib import Path

# ── Required ───────────────────────────────────────────────────────────────
try:
    import fitz  # PyMuPDF
except ImportError:
    sys.exit("Missing dependency.\nRun: python -m pip install pymupdf")

try:
    from PIL import Image, ImageTk, ImageDraw, ImageFont
except ImportError:
    sys.exit("Missing dependency.\nRun: python -m pip install Pillow")

# ── Optional ───────────────────────────────────────────────────────────────
try:
    import docx as _docx_lib
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation as _Pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ── File type registry ─────────────────────────────────────────────────────
SUPPORTED_EXTENSIONS = {
    '.pdf':  'pdf',
    '.jpg':  'image', '.jpeg': 'image', '.png':  'image',
    '.gif':  'image', '.bmp':  'image', '.tiff': 'image', '.tif': 'image',
    '.docx': 'word',
    '.xlsx': 'excel', '.xls': 'excel',
    '.pptx': 'pptx',
}

TYPE_ICON = {
    'pdf':   '📄',
    'image': '🖼',
    'word':  '📝',
    'excel': '📊',
    'pptx':  '📋',
}

BASE_IMG_WIDTH = 860


# ── Utilities ──────────────────────────────────────────────────────────────

def find_files(folder: Path) -> list:
    files = []
    try:
        for entry in os.scandir(folder):
            if entry.is_file():
                ext = Path(entry.name).suffix.lower()
                if ext in SUPPORTED_EXTENSIONS:
                    files.append(Path(entry.path))
    except PermissionError:
        pass
    return sorted(files, key=lambda p: p.name.lower())


def notes_path(folder: Path) -> Path:
    return folder / ".file_review_notes.json"


def load_notes(folder: Path) -> dict:
    p = notes_path(folder)
    if p.exists():
        try:
            with open(p, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}


def save_notes(folder: Path, data: dict) -> None:
    with open(notes_path(folder), "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)


# ── Font helpers ───────────────────────────────────────────────────────────

def get_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    bold_names = [
        "arialbd.ttf", "Arial Bold.ttf",
        "C:/Windows/Fonts/arialbd.ttf",
        "calibrib.ttf", "C:/Windows/Fonts/calibrib.ttf",
    ]
    normal_names = [
        "arial.ttf", "Arial.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "calibri.ttf", "C:/Windows/Fonts/calibri.ttf",
        "DejaVuSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for name in (bold_names if bold else normal_names):
        try:
            return ImageFont.truetype(name, size)
        except Exception:
            pass
    return ImageFont.load_default()


# ── Text-to-image renderer ─────────────────────────────────────────────────

class LineBuffer:
    def __init__(self, width: int = BASE_IMG_WIDTH, margin: int = 28):
        self.width   = width
        self.margin  = margin
        self.entries = []

    def add(self, text: str, font,
            color: str = "#1a1a1a", x_indent: int = 0, gap_before: int = 0):
        self.entries.append((text, font, color, x_indent, gap_before))

    def render(self, warn_header: str = None) -> Image.Image:
        dummy = Image.new("RGB", (1, 1))
        d     = ImageDraw.Draw(dummy)

        total_h = self.margin
        if warn_header:
            total_h += 54
        for text, font, _, _, gap_before in self.entries:
            total_h += gap_before
            bbox     = d.textbbox((0, 0), text, font=font)
            total_h += (bbox[3] - bbox[1]) + 4
        total_h = max(total_h + self.margin, 200)

        img  = Image.new("RGB", (self.width, total_h), "#ffffff")
        draw = ImageDraw.Draw(img)
        y    = self.margin

        if warn_header:
            draw.rectangle([(0, 0), (self.width, 50)],  fill="#fff3cd")
            draw.rectangle([(0, 50), (self.width, 52)], fill="#e6a817")
            draw.text((12, 15), warn_header,
                      font=get_font(11, bold=True), fill="#7d4e00")
            y = 62

        for text, font, color, x_indent, gap_before in self.entries:
            y += gap_before
            draw.text((self.margin + x_indent, y), text, font=font, fill=color)
            bbox = draw.textbbox((0, 0), text, font=font)
            y   += (bbox[3] - bbox[1]) + 4

        return img


# ── Per-type renderers ─────────────────────────────────────────────────────

def render_word(path: Path) -> tuple:
    if not HAS_DOCX:
        return _missing_lib_image("python-docx",
                                  "python -m pip install python-docx"), 1
    doc        = _docx_lib.Document(str(path))
    buf        = LineBuffer()
    f_h1       = get_font(20, bold=True)
    f_h2       = get_font(16, bold=True)
    f_body     = get_font(13)
    wrap_chars = 92
    for para in doc.paragraphs:
        text  = para.text.strip()
        style = para.style.name if para.style else ""
        if not text:
            buf.add("", f_body, gap_before=4)
            continue
        if "Heading 1" in style:
            buf.add(text, f_h1, color="#1a3a5c", gap_before=14)
        elif "Heading 2" in style or "Heading 3" in style:
            buf.add(text, f_h2, color="#2c5282", gap_before=10)
        else:
            for line in textwrap.wrap(text, wrap_chars) or [text]:
                buf.add(line, f_body, gap_before=2)
    warn = (f"⚠  Approximate text view — not a native Word render.  "
            f"Open '{path.name}' in Microsoft Word for the true layout.")
    return buf.render(warn_header=warn), 1


def render_excel(path: Path, sheet_idx: int = 0) -> tuple:
    if not HAS_XLSX:
        return _missing_lib_image("openpyxl",
                                  "python -m pip install openpyxl"), 1
    wb          = openpyxl.load_workbook(str(path), read_only=True,
                                         data_only=True)
    sheet_names = wb.sheetnames
    n_sheets    = len(sheet_names)
    sheet_idx   = max(0, min(sheet_idx, n_sheets - 1))
    ws          = wb[sheet_names[sheet_idx]]
    MAX_ROWS, MAX_COLS = 300, 20
    rows = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i >= MAX_ROWS:
            break
        rows.append([str(c) if c is not None else ""
                     for c in row[:MAX_COLS]])
    wb.close()
    if not rows:
        rows = [["(empty sheet)"]]
    n_cols = max(len(r) for r in rows)
    for r in rows:
        while len(r) < n_cols:
            r.append("")
    f_head  = get_font(11, bold=True)
    f_cell  = get_font(11)
    COL_PAD = 10
    ROW_H   = 22
    MAX_CW  = 160
    col_widths = []
    for c in range(n_cols):
        max_chars = max(len(rows[r][c]) for r in range(len(rows)))
        w = min(max_chars * 7 + COL_PAD * 2, MAX_CW)
        col_widths.append(max(w, 44))
    total_w = max(sum(col_widths) + 2, BASE_IMG_WIDTH)
    total_h = max(len(rows) * ROW_H + 56 + 2, 200)
    img  = Image.new("RGB", (total_w, total_h), "#ffffff")
    draw = ImageDraw.Draw(img)
    draw.rectangle([(0, 0), (total_w, 50)],  fill="#fff3cd")
    draw.rectangle([(0, 50), (total_w, 52)], fill="#e6a817")
    warn = (f"⚠  Approximate table view — not a native Excel render.  "
            f"Open '{path.name}' in Excel for the true layout.  "
            f"Sheet: {sheet_names[sheet_idx]}  ({sheet_idx + 1}/{n_sheets})")
    draw.text((12, 15), warn, font=get_font(11, bold=True), fill="#7d4e00")
    y0 = 54
    for ri, row in enumerate(rows):
        x        = 1
        bg       = "#dce8fc" if ri == 0 else \
                   ("#f4f4f4" if ri % 2 == 0 else "#ffffff")
        row_font = f_head if ri == 0 else f_cell
        for ci, cell in enumerate(row):
            cw = col_widths[ci]
            draw.rectangle(
                [(x, y0 + ri * ROW_H),
                 (x + cw - 1, y0 + (ri + 1) * ROW_H - 1)],
                fill=bg, outline="#d0d0d0")
            display = cell
            while (display and
                   draw.textlength(display, font=row_font) > cw - COL_PAD * 2):
                display = display[:-1]
            if display != cell:
                display = display[:-1] + "…"
            draw.text((x + COL_PAD, y0 + ri * ROW_H + 4),
                      display, font=row_font, fill="#1a1a1a")
            x += cw
    return img, n_sheets


def render_pptx(path: Path, slide_idx: int = 0) -> tuple:
    if not HAS_PPTX:
        return _missing_lib_image("python-pptx",
                                  "python -m pip install python-pptx"), 1
    prs       = _Pptx(str(path))
    n_slides  = len(prs.slides)
    slide_idx = max(0, min(slide_idx, n_slides - 1))
    slide     = prs.slides[slide_idx]
    buf     = LineBuffer()
    f_title = get_font(22, bold=True)
    f_body  = get_font(13)
    f_small = get_font(11)
    first   = True
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = para.level
            if first:
                buf.add(text, f_title, color="#1a3a5c", gap_before=8)
                first = False
            elif level == 0:
                for line in textwrap.wrap(text, 90) or [text]:
                    buf.add(line, f_body, gap_before=8)
            else:
                indent = level * 16
                for line in textwrap.wrap(text, 90 - level * 4) or [text]:
                    buf.add("• " + line, f_small,
                            x_indent=indent, gap_before=3)
    if first:
        buf.add("(No text content on this slide)",
                f_body, color="#888888", gap_before=20)
    warn = (f"⚠  Text-only view — not a native PowerPoint render.  "
            f"Open '{path.name}' in PowerPoint for the true layout.  "
            f"Slide {slide_idx + 1} of {n_slides}.")
    return buf.render(warn_header=warn), n_slides


def _missing_lib_image(lib_name: str, install_cmd: str) -> Image.Image:
    buf = LineBuffer()
    f   = get_font(13)
    buf.add("Cannot display this file type.",
            f, color="#c0392b", gap_before=30)
    buf.add(f"Required library not installed:  {lib_name}",
            f, color="#333333", gap_before=12)
    buf.add("Install it by running in PowerShell:",
            f, color="#333333", gap_before=8)
    buf.add(f"    {install_cmd}",
            get_font(13), color="#0055cc", gap_before=4)
    buf.add("Then restart the browser.",
            f, color="#333333", gap_before=8)
    return buf.render()


# ── Main Application ───────────────────────────────────────────────────────

class DocBrowser(tk.Tk):
    ZOOM_STEP  = 0.2
    ZOOM_MIN   = 0.3
    ZOOM_MAX   = 4.0
    PDF_BASE_H = 860

    def __init__(self):
        super().__init__()
        self.title("Advanced File Review")
        self.configure(bg="#1e1e2e")
        self.resizable(True, True)

        self.folder   = None
        self.files    = []
        self.notes    = {}
        self.idx      = 0
        self.page_no  = 0
        self.zoom     = 1.0
        self._ftype   = None
        self._doc     = None
        self._img_ref = None

        self._build_ui()
        self.geometry("1500x880")
        self.minsize(900, 600)

    # ══════════════════════════════════════════════════════════════════════
    # UI Construction
    # ══════════════════════════════════════════════════════════════════════

    def _build_ui(self):
        self._build_folder_bar()
        self._build_main_area()

    # ── Top bar: folder selection ──────────────────────────────────────────

    def _build_folder_bar(self):
        bar = tk.Frame(self, bg="#181825", pady=8)
        bar.pack(fill="x")

        tk.Label(bar, text="Folder:", bg="#181825", fg="#a6adc8",
                 font=("Helvetica", 10, "bold")).pack(side="left", padx=(14, 6))

        self.var_folder = tk.StringVar()
        self.ent_folder = tk.Entry(
            bar, textvariable=self.var_folder,
            bg="#313244", fg="#cdd6f4",
            insertbackground="#cdd6f4", relief="flat",
            font=("Helvetica", 10), width=70)
        self.ent_folder.pack(side="left", padx=(0, 6))
        # Allow pressing Enter to load
        self.ent_folder.bind("<Return>", lambda _: self._load_folder())

        _btn = dict(relief="flat", font=("Helvetica", 9, "bold"),
                    padx=10, pady=4, cursor="hand2")

        tk.Button(bar, text="Browse…", command=self._browse_folder,
                  bg="#313244", fg="#cdd6f4",
                  activebackground="#45475a", activeforeground="#cdd6f4",
                  **_btn).pack(side="left", padx=(0, 4))

        tk.Button(bar, text="⟳  Load Files", command=self._load_folder,
                  bg="#89b4fa", fg="#1e1e2e",
                  activebackground="#74c7ec", activeforeground="#1e1e2e",
                  **_btn).pack(side="left")

        self.lbl_count = tk.Label(bar, text="", bg="#181825",
                                  fg="#6c7086", font=("Helvetica", 9))
        self.lbl_count.pack(side="right", padx=14)

        self.lbl_file = tk.Label(bar, text="No folder loaded",
                                 bg="#181825", fg="#89b4fa",
                                 font=("Helvetica", 10, "bold"), anchor="w")
        self.lbl_file.pack(side="right", padx=(0, 20))

    # ── Three-pane main area ───────────────────────────────────────────────

    def _build_main_area(self):
        pw = tk.PanedWindow(self, orient="horizontal", bg="#1e1e2e",
                            sashwidth=5, sashpad=0, sashrelief="flat")
        pw.pack(fill="both", expand=True, padx=6, pady=(0, 6))

        # Left: file list
        left = tk.Frame(pw, bg="#1e1e2e")
        pw.add(left, minsize=160, stretch="never")

        # Center: document viewer
        center = tk.Frame(pw, bg="#181825")
        pw.add(center, minsize=400, stretch="always")

        # Right: notes
        right = tk.Frame(pw, bg="#1e1e2e")
        pw.add(right, minsize=260, stretch="never")

        pw.paneconfig(left,   width=220)
        pw.paneconfig(center, width=940)
        pw.paneconfig(right,  width=300)

        self._build_file_list(left)
        self._build_viewer(center)
        self._build_notes_panel(right)

    # ── File list panel (left) ─────────────────────────────────────────────

    def _build_file_list(self, parent):
        hdr = tk.Frame(parent, bg="#1e1e2e")
        hdr.pack(fill="x", padx=6, pady=(8, 4))
        tk.Label(hdr, text="FILES", bg="#1e1e2e", fg="#a6adc8",
                 font=("Helvetica", 9, "bold")).pack(side="left")

        # Listbox + scrollbar
        lf = tk.Frame(parent, bg="#181825")
        lf.pack(fill="both", expand=True, padx=6, pady=(0, 6))

        vsb = ttk.Scrollbar(lf, orient="vertical")
        self.file_listbox = tk.Listbox(
            lf,
            bg="#181825", fg="#cdd6f4",
            selectbackground="#45475a", selectforeground="#89b4fa",
            font=("Helvetica", 9), relief="flat",
            activestyle="none",
            cursor="hand2",
            yscrollcommand=vsb.set,
        )
        vsb.config(command=self.file_listbox.yview)
        vsb.pack(side="right", fill="y")
        self.file_listbox.pack(fill="both", expand=True)

        self.file_listbox.bind("<<ListboxSelect>>", self._on_list_select)

    def _refresh_file_list(self):
        """Rebuild the listbox contents and highlight the current file."""
        self.file_listbox.delete(0, "end")
        for fp in self.files:
            ext   = fp.suffix.lower()
            ftype = SUPPORTED_EXTENSIONS.get(ext, "")
            icon  = TYPE_ICON.get(ftype, "📎")
            self.file_listbox.insert("end", f" {icon}  {fp.name}")

        self._highlight_list_selection()

    def _highlight_list_selection(self):
        """Keep the listbox selection in sync with self.idx."""
        if not self.files:
            return
        self.file_listbox.selection_clear(0, "end")
        self.file_listbox.selection_set(self.idx)
        self.file_listbox.see(self.idx)

    def _on_list_select(self, event):
        sel = self.file_listbox.curselection()
        if sel and sel[0] != self.idx:
            self._open_file(sel[0])

    # ── Viewer panel (center) ──────────────────────────────────────────────

    def _build_viewer(self, parent):
        # Sub-bar: prev/next + page nav + zoom
        sub = tk.Frame(parent, bg="#181825")
        sub.pack(fill="x", pady=(4, 0))

        sm = dict(bg="#313244", fg="#cdd6f4", relief="flat",
                  font=("Helvetica", 9), padx=6, pady=2,
                  cursor="hand2", activebackground="#45475a",
                  activeforeground="#cdd6f4")

        self.btn_prev = tk.Button(sub, text="◀ Prev",
                                  command=self._prev, **sm)
        self.btn_prev.pack(side="left", padx=(4, 2))

        self.btn_next = tk.Button(sub, text="Next ▶",
                                  command=self._next, **sm)
        self.btn_next.pack(side="left", padx=(0, 10))

        # Page nav
        self.btn_pg_prev = tk.Button(sub, text="← page",
                                     command=self._page_prev, **sm)
        self.btn_pg_prev.pack(side="left", padx=(0, 0))

        self.lbl_page = tk.Label(sub, text="", bg="#181825",
                                 fg="#6c7086", font=("Helvetica", 9))
        self.lbl_page.pack(side="left", padx=2)

        self.btn_pg_next = tk.Button(sub, text="page →",
                                     command=self._page_next, **sm)
        self.btn_pg_next.pack(side="left", padx=(0, 14))

        # Zoom
        tk.Label(sub, text="Zoom:", bg="#181825", fg="#a6adc8",
                 font=("Helvetica", 9)).pack(side="left", padx=(4, 2))

        self.btn_zoom_out = tk.Button(sub, text=" − ",
                                      command=self._zoom_out, **sm)
        self.btn_zoom_out.pack(side="left", padx=1)

        self.lbl_zoom = tk.Label(sub, text="100%", bg="#181825",
                                 fg="#cdd6f4", font=("Helvetica", 9), width=5)
        self.lbl_zoom.pack(side="left")

        self.btn_zoom_in = tk.Button(sub, text=" + ",
                                     command=self._zoom_in, **sm)
        self.btn_zoom_in.pack(side="left", padx=1)

        tk.Button(sub, text="Reset", command=self._zoom_reset,
                  **sm).pack(side="left", padx=(4, 0))

        # Scrollable canvas
        cf = tk.Frame(parent, bg="#181825")
        cf.pack(fill="both", expand=True)

        self.canvas = tk.Canvas(cf, bg="#181825", highlightthickness=0)
        vsb = ttk.Scrollbar(cf, orient="vertical",   command=self.canvas.yview)
        hsb = ttk.Scrollbar(cf, orient="horizontal", command=self.canvas.xview)
        self.canvas.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)
        vsb.pack(side="right",  fill="y")
        hsb.pack(side="bottom", fill="x")
        self.canvas.pack(fill="both", expand=True)

        self.canvas.bind("<MouseWheel>", self._on_mousewheel)
        self.canvas.bind("<Button-4>",   self._on_mousewheel)
        self.canvas.bind("<Button-5>",   self._on_mousewheel)

    # ── Notes panel (right) ────────────────────────────────────────────────

    def _build_notes_panel(self, parent):
        pad = dict(padx=14, pady=6)
        lbl = dict(bg="#1e1e2e", fg="#a6adc8",
                   font=("Helvetica", 10, "bold"), anchor="w")
        ent = dict(bg="#313244", fg="#cdd6f4",
                   insertbackground="#cdd6f4", relief="flat",
                   font=("Helvetica", 10))

        tk.Label(parent, text="DOCUMENT NOTES", **lbl).pack(
            fill="x", padx=14, pady=(16, 2))

        tk.Label(parent, text="Document Date", **lbl).pack(fill="x", **pad)
        self.ent_date = tk.Entry(parent, **ent)
        self.ent_date.pack(fill="x", padx=14)

        tk.Label(parent, text="Notes", **lbl).pack(fill="x", **pad)
        self.txt_notes = tk.Text(
            parent, bg="#313244", fg="#cdd6f4",
            insertbackground="#cdd6f4", relief="flat",
            font=("Helvetica", 10), wrap="word", height=14)
        self.txt_notes.pack(fill="both", expand=True, padx=14)

        tk.Button(parent, text="💾  Save Notes",
                  bg="#89b4fa", fg="#1e1e2e", relief="flat",
                  font=("Helvetica", 10, "bold"), pady=6, cursor="hand2",
                  activebackground="#74c7ec", activeforeground="#1e1e2e",
                  command=self._save_current_notes).pack(
            fill="x", padx=14, pady=(10, 4))

        tk.Frame(parent, bg="#313244", height=1).pack(
            fill="x", padx=14, pady=8)

        tk.Label(parent, text="EXPORT ALL NOTES", **lbl).pack(
            fill="x", padx=14, pady=(2, 6))

        tk.Button(parent, text="📄  Export as CSV",
                  bg="#a6e3a1", fg="#1e1e2e", relief="flat",
                  font=("Helvetica", 10, "bold"), pady=6, cursor="hand2",
                  activebackground="#94e2d5", activeforeground="#1e1e2e",
                  command=lambda: self._export("csv")).pack(
            fill="x", padx=14, pady=2)

        tk.Button(parent, text="📝  Export as TXT",
                  bg="#f9e2af", fg="#1e1e2e", relief="flat",
                  font=("Helvetica", 10, "bold"), pady=6, cursor="hand2",
                  activebackground="#fab387", activeforeground="#1e1e2e",
                  command=lambda: self._export("txt")).pack(
            fill="x", padx=14, pady=2)

        self.lbl_status = tk.Label(parent, text="", bg="#1e1e2e",
                                   fg="#a6e3a1", font=("Helvetica", 9),
                                   wraplength=240)
        self.lbl_status.pack(fill="x", padx=14, pady=(6, 12))

    # ══════════════════════════════════════════════════════════════════════
    # Folder loading
    # ══════════════════════════════════════════════════════════════════════

    def _browse_folder(self):
        d = filedialog.askdirectory(title="Select folder to review")
        if d:
            self.var_folder.set(d)
            self._load_folder()

    def _load_folder(self):
        raw = self.var_folder.get().strip()
        if not raw:
            messagebox.showerror("No folder", "Please enter or browse to a folder.")
            return
        folder = Path(raw)
        if not folder.is_dir():
            messagebox.showerror("Not found",
                                 f"Could not find folder:\n{raw}\n\n"
                                 "Check the path and try again.")
            return

        # Save any unsaved notes before switching folders
        if self.folder and self.files:
            self._save_current_notes()

        self.folder = folder
        self.files  = find_files(folder)
        self.notes  = load_notes(folder)
        self.idx    = 0
        self.page_no = 0

        n = len(self.files)
        self.lbl_count.config(
            text=f"{n} file{'s' if n != 1 else ''} found")

        self._refresh_file_list()

        if self.files:
            self._open_file(0)
        else:
            self._show_empty()

    # ══════════════════════════════════════════════════════════════════════
    # File opening & rendering
    # ══════════════════════════════════════════════════════════════════════

    def _open_file(self, idx: int):
        self._save_current_notes()
        self.idx     = idx
        self.page_no = 0

        if self._doc:
            self._doc.close()
            self._doc = None

        path        = self.files[idx]
        ext         = path.suffix.lower()
        self._ftype = SUPPORTED_EXTENSIONS.get(ext, "unknown")

        self.lbl_file.config(text=path.name)
        self.lbl_count.config(
            text=f"{idx + 1} / {len(self.files)}")

        if self._ftype == "pdf":
            self._doc = fitz.open(str(path))

        self._highlight_list_selection()
        self._render_current()
        self._load_notes_fields(path.name)
        self._update_nav_buttons()

    def _render_current(self):
        path  = self.files[self.idx]
        ftype = self._ftype

        if ftype == "pdf":
            self._render_pdf(path)
        elif ftype == "image":
            self._render_image(path)
        elif ftype == "word":
            img, n = render_word(path)
            self._display_pil(img, page_label=None, total=n)
        elif ftype == "excel":
            img, n = render_excel(path, self.page_no)
            self._display_pil(img,
                              page_label=f"Sheet {self.page_no + 1} of {n}",
                              total=n)
        elif ftype == "pptx":
            img, n = render_pptx(path, self.page_no)
            self._display_pil(img,
                              page_label=f"Slide {self.page_no + 1} of {n}",
                              total=n)
        else:
            self._display_pil(
                _missing_lib_image("(unknown type)", ""), None, 1)

    def _render_pdf(self, path: Path):
        page   = self._doc[self.page_no]
        zoom   = (self.PDF_BASE_H / page.rect.height) * self.zoom
        matrix = fitz.Matrix(zoom, zoom)
        pix    = page.get_pixmap(matrix=matrix, alpha=False)
        img    = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        n      = len(self._doc)
        self._display_pil(img,
                          page_label=f"Page {self.page_no + 1} of {n}",
                          total=n)

    def _render_image(self, path: Path):
        img = Image.open(str(path)).convert("RGB")
        w   = int(img.width  * self.zoom)
        h   = int(img.height * self.zoom)
        img = img.resize((w, h), Image.LANCZOS)
        self._display_pil(img, page_label=None, total=1)

    def _display_pil(self, img: Image.Image, page_label, total: int):
        if self._ftype not in ("pdf", "image") and self.zoom != 1.0:
            img = img.resize(
                (int(img.width * self.zoom),
                 int(img.height * self.zoom)),
                Image.LANCZOS)

        self._img_ref = ImageTk.PhotoImage(img)
        self.canvas.delete("all")
        self.canvas.create_image(4, 4, anchor="nw", image=self._img_ref)
        self.canvas.config(scrollregion=(0, 0,
                                         img.width  + 8,
                                         img.height + 8))
        if page_label:
            self.lbl_page.config(text=f"  {page_label}  ")
            self.btn_pg_prev.config(
                state="normal" if self.page_no > 0 else "disabled")
            self.btn_pg_next.config(
                state="normal" if self.page_no < total - 1 else "disabled")
        else:
            self.lbl_page.config(text="")
            self.btn_pg_prev.config(state="disabled")
            self.btn_pg_next.config(state="disabled")

    # ══════════════════════════════════════════════════════════════════════
    # Navigation
    # ══════════════════════════════════════════════════════════════════════

    def _prev(self):
        if self.idx > 0:
            self._open_file(self.idx - 1)

    def _next(self):
        if self.idx < len(self.files) - 1:
            self._open_file(self.idx + 1)

    def _page_prev(self):
        if self.page_no > 0:
            self.page_no -= 1
            self._render_current()

    def _page_next(self):
        self.page_no += 1
        self._render_current()

    def _update_nav_buttons(self):
        self.btn_prev.config(
            state="normal" if self.idx > 0 else "disabled")
        self.btn_next.config(
            state="normal" if self.idx < len(self.files) - 1 else "disabled")

    # ══════════════════════════════════════════════════════════════════════
    # Zoom
    # ══════════════════════════════════════════════════════════════════════

    def _zoom_in(self):
        self.zoom = min(round(self.zoom + self.ZOOM_STEP, 1), self.ZOOM_MAX)
        self._apply_zoom()

    def _zoom_out(self):
        self.zoom = max(round(self.zoom - self.ZOOM_STEP, 1), self.ZOOM_MIN)
        self._apply_zoom()

    def _zoom_reset(self):
        self.zoom = 1.0
        self._apply_zoom()

    def _apply_zoom(self):
        self.lbl_zoom.config(text=f"{int(self.zoom * 100)}%")
        self._render_current()

    def _on_mousewheel(self, event):
        if event.num == 4:
            self.canvas.yview_scroll(-1, "units")
        elif event.num == 5:
            self.canvas.yview_scroll(1, "units")
        else:
            self.canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

    def _show_empty(self):
        self.lbl_file.config(text="No supported files found.")
        self.canvas.delete("all")
        self.canvas.create_text(
            300, 200,
            text="No supported files found in this folder.\n\n"
                 "Supported types: PDF, JPG/PNG/TIFF, DOCX, XLSX, PPTX",
            fill="#585b70", font=("Helvetica", 13), justify="center")

    # ══════════════════════════════════════════════════════════════════════
    # Notes
    # ══════════════════════════════════════════════════════════════════════

    def _load_notes_fields(self, filename: str):
        entry = self.notes.get(filename, {})
        self.ent_date.delete(0, "end")
        self.ent_date.insert(0, entry.get("date", ""))
        self.txt_notes.delete("1.0", "end")
        self.txt_notes.insert("1.0", entry.get("notes", ""))

    def _save_current_notes(self):
        if not self.files or not self.folder:
            return
        filename = self.files[self.idx].name
        self.notes[filename] = {
            "date":  self.ent_date.get().strip(),
            "notes": self.txt_notes.get("1.0", "end").strip(),
        }
        save_notes(self.folder, self.notes)

    # ══════════════════════════════════════════════════════════════════════
    # Export
    # ══════════════════════════════════════════════════════════════════════

    def _export(self, fmt: str):
        if not self.folder or not self.files:
            messagebox.showwarning("Nothing to export",
                                   "Load a folder first.")
            return
        self._save_current_notes()
        stamp   = datetime.now().strftime("%Y%m%d_%H%M%S")
        default = f"doc_notes_{stamp}"

        if fmt == "csv":
            path = filedialog.asksaveasfilename(
                defaultextension=".csv",
                filetypes=[("CSV files", "*.csv"), ("All files", "*")],
                initialfile=default + ".csv",
                title="Export notes as CSV")
        else:
            path = filedialog.asksaveasfilename(
                defaultextension=".txt",
                filetypes=[("Text files", "*.txt"), ("All files", "*")],
                initialfile=default + ".txt",
                title="Export notes as TXT")

        if not path:
            return

        try:
            if fmt == "csv":
                self._write_csv(path)
            else:
                self._write_txt(path)
            self.lbl_status.config(
                text=f"✅ Exported:\n{Path(path).name}", fg="#a6e3a1")
        except Exception as e:
            messagebox.showerror("Export failed", str(e))
            self.lbl_status.config(text="❌ Export failed.", fg="#f38ba8")

    def _write_csv(self, path: str):
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.DictWriter(
                f, fieldnames=["filename", "type", "date", "notes"])
            writer.writeheader()
            for fp in self.files:
                entry = self.notes.get(fp.name, {})
                ftype = SUPPORTED_EXTENSIONS.get(fp.suffix.lower(), "unknown")
                writer.writerow({
                    "filename": fp.name,
                    "type":     ftype,
                    "date":     entry.get("date", ""),
                    "notes":    entry.get("notes", ""),
                })

    def _write_txt(self, path: str):
        with open(path, "w", encoding="utf-8") as f:
            f.write("DOCUMENT REVIEW — NOTES EXPORT\n")
            f.write(f"Folder:    {self.folder}\n")
            f.write(f"Generated: "
                    f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write("=" * 60 + "\n\n")
            for fp in self.files:
                entry = self.notes.get(fp.name, {})
                ftype = SUPPORTED_EXTENSIONS.get(fp.suffix.lower(), "unknown")
                f.write(f"FILE:  {fp.name}\n")
                f.write(f"TYPE:  {ftype}\n")
                f.write(f"DATE:  {entry.get('date', '(none)')}\n")
                f.write(f"NOTES:\n{entry.get('notes', '(none)')}\n")
                f.write("-" * 60 + "\n\n")

    # ══════════════════════════════════════════════════════════════════════
    # Lifecycle
    # ══════════════════════════════════════════════════════════════════════

    def destroy(self):
        self._save_current_notes()
        if self._doc:
            self._doc.close()
        super().destroy()


# ── Entry point ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    app = DocBrowser()
    app.mainloop()
